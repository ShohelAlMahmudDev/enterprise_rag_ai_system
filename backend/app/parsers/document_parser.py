import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.config import settings
from app.services.ingestion_service import ImageIngestionService
from app.utils.ocr_utils import extract_ocr_text_from_pdf_page

logger = logging.getLogger(__name__)

SUPPORTED_IMAGE_EXTENSIONS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".tif",
    ".tiff",
    ".gif",
}


@dataclass(slots=True)
class ParsedSection:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(slots=True)
class ParsedDocument:
    text: str
    sections: list[ParsedSection]
    file_type: str


def parse_document(file_path: str) -> ParsedDocument:
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".txt":
        return _parse_txt(path)
    if ext == ".pdf":
        return _parse_pdf(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext == ".xlsx":
        return _parse_xlsx(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    if ext in SUPPORTED_IMAGE_EXTENSIONS:
        return _parse_image(path)

    raise ValueError(f"Unsupported parser for extension: {ext}")


def _normalize_text(text: str) -> str:
    if not text:
        return ""

    lines = [line.strip() for line in text.splitlines()]
    filtered = [line for line in lines if line]
    return "\n".join(filtered).strip()


def _build_parsed_document(
    *,
    sections: list[ParsedSection],
    file_type: str,
) -> ParsedDocument:
    full_text = "\n\n".join(section.text for section in sections if section.text).strip()
    return ParsedDocument(
        text=full_text,
        sections=sections,
        file_type=file_type,
    )


def _parse_txt(path: Path) -> ParsedDocument:
    text = path.read_text(encoding="utf-8", errors="ignore")
    text = _normalize_text(text)

    sections: list[ParsedSection] = []
    if text:
        sections.append(
            ParsedSection(
                text=text,
                metadata={
                    "type": "text",
                    "source_file": path.name,
                    "file_type": ".txt",
                },
            )
        )

    return _build_parsed_document(sections=sections, file_type=".txt")


def _parse_pdf(path: Path) -> ParsedDocument:
    reader = PdfReader(str(path))
    sections: list[ParsedSection] = []

    for page_number, page in enumerate(reader.pages, start=1):
        extracted = page.extract_text() or ""
        cleaned = _normalize_text(extracted)
        used_ocr = False

        if settings.OCR_ENABLED and len(cleaned) < settings.PDF_OCR_MIN_TEXT_LENGTH:
            try:
                ocr_text = extract_ocr_text_from_pdf_page(path, page_number)
                ocr_cleaned = _normalize_text(ocr_text)

                if len(ocr_cleaned) > len(cleaned):
                    cleaned = ocr_cleaned
                    used_ocr = True
            except Exception as exc:
                logger.warning(
                    "OCR failed for PDF page %s in %s: %s",
                    page_number,
                    path.name,
                    exc,
                )
                used_ocr = False

        if not cleaned:
            continue

        prefix = f"Page {page_number}"
        if used_ocr:
            prefix += " (OCR)"

        page_text = f"{prefix}\n{cleaned}"
        sections.append(
            ParsedSection(
                text=page_text,
                metadata={
                    "type": "pdf_page",
                    "page": page_number,
                    "ocr_used": used_ocr,
                    "source_file": path.name,
                    "file_type": ".pdf",
                },
            )
        )

    return _build_parsed_document(sections=sections, file_type=".pdf")


def _parse_docx(path: Path) -> ParsedDocument:
    doc = DocxDocument(str(path))
    sections: list[ParsedSection] = []

    current_heading = "Document"
    buffer: list[str] = []

    def flush_buffer() -> None:
        nonlocal buffer

        cleaned = _normalize_text("\n".join(buffer))
        if cleaned:
            section_text = f"Section: {current_heading}\n{cleaned}"
            sections.append(
                ParsedSection(
                    text=section_text,
                    metadata={
                        "type": "docx_section",
                        "heading": current_heading,
                        "source_file": path.name,
                        "file_type": ".docx",
                    },
                )
            )
        buffer = []

    for paragraph in doc.paragraphs:
        text = (paragraph.text or "").strip()
        if not text:
            continue

        style_name = paragraph.style.name.lower() if paragraph.style and paragraph.style.name else ""

        if "heading" in style_name:
            flush_buffer()
            current_heading = text
        else:
            buffer.append(text)

    flush_buffer()

    for table_index, table in enumerate(doc.tables, start=1):
        rows: list[str] = []

        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if values:
                rows.append(" | ".join(values))

        cleaned_table = _normalize_text("\n".join(rows))
        if cleaned_table:
            sections.append(
                ParsedSection(
                    text=f"Table {table_index}\n{cleaned_table}",
                    metadata={
                        "type": "docx_table",
                        "table_index": table_index,
                        "source_file": path.name,
                        "file_type": ".docx",
                    },
                )
            )

    return _build_parsed_document(sections=sections, file_type=".docx")


def _parse_xlsx(path: Path) -> ParsedDocument:
    wb = load_workbook(str(path), data_only=True)
    sections: list[ParsedSection] = []

    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        header: list[str] | None = None
        data_rows: list[tuple[Any, ...]] = []

        for row in rows:
            normalized = [str(v).strip() if v is not None else "" for v in row]
            if any(normalized):
                if header is None:
                    header = normalized
                else:
                    data_rows.append(row)

        if not header:
            continue

        summary_lines = [f"Sheet: {ws.title}"]
        summary_lines.append("Columns: " + ", ".join(col for col in header if col))
        sections.append(
            ParsedSection(
                text=_normalize_text("\n".join(summary_lines)),
                metadata={
                    "type": "xlsx_sheet_summary",
                    "sheet": ws.title,
                    "source_file": path.name,
                    "file_type": ".xlsx",
                },
            )
        )

        for row_index, row in enumerate(data_rows, start=2):
            values = [str(v).strip() if v is not None else "" for v in row]
            pairs: list[str] = []

            for col_name, value in zip(header, values):
                normalized_col_name = (col_name or "").strip()
                normalized_value = value.strip()
                if normalized_col_name and normalized_value:
                    pairs.append(f"{normalized_col_name}: {normalized_value}")

            if pairs:
                row_text = f"Sheet: {ws.title}\nRow: {row_index}\n" + "\n".join(pairs)
                sections.append(
                    ParsedSection(
                        text=_normalize_text(row_text),
                        metadata={
                            "type": "xlsx_row",
                            "sheet": ws.title,
                            "row": row_index,
                            "source_file": path.name,
                            "file_type": ".xlsx",
                        },
                    )
                )

    return _build_parsed_document(sections=sections, file_type=".xlsx")


def _parse_pptx(path: Path) -> ParsedDocument:
    prs = Presentation(str(path))
    sections: list[ParsedSection] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []

        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text and text != title:
                    body_parts.append(text)

        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide:
                note_parts: list[str] = []
                for shape in slide.notes_slide.shapes:
                    if hasattr(shape, "text"):
                        note_text = (shape.text or "").strip()
                        if note_text:
                            note_parts.append(note_text)
                notes_text = "\n".join(note_parts).strip()
        except Exception as exc:
            logger.warning(
                "Failed to read notes for slide %s in %s: %s",
                slide_index,
                path.name,
                exc,
            )
            notes_text = ""

        parts = [f"Slide {slide_index}"]
        if title:
            parts.append(f"Title: {title}")
        if body_parts:
            parts.append("Content:")
            parts.extend(body_parts)
        if notes_text:
            parts.append("Notes:")
            parts.append(notes_text)

        slide_text = _normalize_text("\n".join(parts))
        if slide_text:
            sections.append(
                ParsedSection(
                    text=slide_text,
                    metadata={
                        "type": "pptx_slide",
                        "slide": slide_index,
                        "title": title or None,
                        "source_file": path.name,
                        "file_type": ".pptx",
                    },
                )
            )

    return _build_parsed_document(sections=sections, file_type=".pptx")


def _parse_image(path: Path) -> ParsedDocument:
    image_type = _infer_image_type_from_name(path.name)
    image_service = ImageIngestionService()

    try:
        chunk = image_service.ingest_image(
            image_path=path,
            logical_name=path.stem,
            filename_hint=path.name,
            image_type=image_type,
            chunk_id=f"{path.stem}_image_1",
        )
    except Exception as exc:
        logger.exception("Failed to parse image %s: %s", path.name, exc)
        raise RuntimeError(f"Failed to parse image file: {path.name}") from exc

    text = _normalize_text(chunk.get("text", ""))
    metadata = dict(chunk.get("metadata", {}))
    metadata.update(
        {
            "type": "image",
            "source_file": path.name,
            "file_type": path.suffix.lower(),
        }
    )

    sections: list[ParsedSection] = []
    if text:
        sections.append(
            ParsedSection(
                text=text,
                metadata=metadata,
            )
        )

    return _build_parsed_document(sections=sections, file_type=path.suffix.lower())


def _infer_image_type_from_name(filename: str) -> str:
    lowered = filename.lower()

    if any(keyword in lowered for keyword in ("diagram", "flow", "state", "sequence", "arch", "network")):
        return "diagram"

    if any(keyword in lowered for keyword in ("screen", "screenshot", "ui", "dashboard")):
        return "screenshot"

    if any(keyword in lowered for keyword in ("scan", "document", "page")):
        return "document"

    if any(keyword in lowered for keyword in ("photo", "img", "image", "picture")):
        return "photo"

    return "unknown"